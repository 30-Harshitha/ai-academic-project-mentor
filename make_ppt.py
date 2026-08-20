import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
DARK_BG = RGBColor(15, 23, 42)      # Dark Slate/Navy #0f172a
LIGHT_BG = RGBColor(248, 250, 252)  # Clean Light #f8fafc
PRIMARY_BLUE = RGBColor(37, 99, 235) # Royal Blue #2563eb
ACCENT_CYAN = RGBColor(6, 182, 212) # Cyan #06b6d4
TEXT_DARK = RGBColor(30, 41, 59)    # Dark Text
TEXT_MUTED = RGBColor(100, 116, 139)# Muted Text
CARD_BG = RGBColor(255, 255, 255)   # White Card
ACCENT_GREEN = RGBColor(16, 185, 129)# Success Green #10b981

def add_header(slide, title_text, category_text='AI ACADEMIC PROJECT MENTOR'):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_CYAN
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)

def create_title_slide(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = 'AI ACADEMIC PROJECT MENTOR'
    p0.font.size = Pt(40)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(255, 255, 255)
    
    p1 = tf.add_paragraph()
    p1.text = 'An Agentic AI Platform for End-to-End Academic Project Lifecycle Management'
    p1.font.size = Pt(20)
    p1.font.color.rgb = ACCENT_CYAN
    p1.space_before = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = 'Program Schedule Timeline: June 29, 2026 – August 20, 2026'
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    p2.space_before = Pt(20)
    
    teamBox = slide.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.333), Inches(2.2))
    ttf = teamBox.text_frame
    ttf.word_wrap = True
    
    tp0 = ttf.paragraphs[0]
    tp0.text = 'PROJECT TEAM MEMBERS & RESPONSIBILITY ALLOCATION:'
    tp0.font.size = Pt(14)
    tp0.font.bold = True
    tp0.font.color.rgb = ACCENT_CYAN
    
    tp1 = ttf.add_paragraph()
    tp1.text = '• Harshitha H S — MAJOR ROLE (Project Lead & Core Multi-Agent Pipeline Architect)\n• Vamshi Krishna — Minor Role (Backend Database & REST API Specialist)\n• Muthumenen M — Minor Role (Frontend UI/UX & Student Onboarding)\n• Lohith Raj — Minor Role (Risk Assessment & AI Mentor Interaction)\n• Harsh — Minor Role (Report Exporter & Faculty Dashboard Developer)'
    tp1.font.size = Pt(13)
    tp1.font.color.rgb = RGBColor(241, 245, 249)
    tp1.space_before = Pt(8)
    return slide

def add_content_slide(prs, title, points, category='PROJECT OVERVIEW'):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, title, category)
    
    for idx, (head, body) in enumerate(points):
        top_pos = Inches(1.4 + (idx % 4) * 1.35)
        left_pos = Inches(0.8)
        width_pos = Inches(11.733)
        height_pos = Inches(1.2)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, width_pos, height_pos)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = RGBColor(226, 232, 240)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.3)
        ctf.margin_top = Inches(0.15)
        
        cp0 = ctf.paragraphs[0]
        cp0.text = head
        cp0.font.size = Pt(16)
        cp0.font.bold = True
        cp0.font.color.rgb = PRIMARY_BLUE
        
        cp1 = ctf.add_paragraph()
        cp1.text = body
        cp1.font.size = Pt(13)
        cp1.font.color.rgb = TEXT_DARK
        cp1.space_before = Pt(4)
    return slide

def add_table_slide(prs, title, headers, rows, category='AGILE MANAGEMENT'):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, title, category)
    
    rows_count = len(rows) + 1
    cols_count = len(headers)
    
    table_shape = slide.shapes.add_table(rows_count, cols_count, Inches(0.8), Inches(1.4), Inches(11.733), Inches(5.5))
    table = table_shape.table
    
    for c_idx, h_text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h_text
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
    for r_idx, row_data in enumerate(rows, start=1):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            if r_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
            else:
                cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = str(val)
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_DARK
            if 'MAJOR' in str(val) or 'PASS' in str(val) or 'Closed' in str(val) or 'Completed' in str(val):
                p.font.bold = True
                if 'MAJOR' in str(val): p.font.color.rgb = PRIMARY_BLUE
                if 'PASS' in str(val) or 'Closed' in str(val) or 'Completed' in str(val): p.font.color.rgb = ACCENT_GREEN
    return slide

# 1. Slide 1: Title
create_title_slide(prs)

# 2. Slide 2: Executive Summary & Project Vision
add_content_slide(prs, 'Executive Summary & Project Vision', [
    ('Bridging the Guidance Gap', 'Develop an intelligent, agentic platform that guides students through the complete academic project lifecycle automatically from a 2-line input.'),
    ('Single-Click Multi-Agent Pipeline', 'Triggers 5 specialized AI agents to generate feasibility ratings, scope boundaries, tech stack recommendations, milestone timelines, and risk mitigations.'),
    ('Conversational Mentoring & On-Demand Reports', 'Provides ongoing interactive mentoring (Nova AI Desk), weekly progress check-in tracking, and dynamic PDF/synopsis documentation exports.'),
    ('Passive Faculty Supervision', 'Empowers teachers with a centralized monitoring dashboard, student team health indicators, and an interactive feedback modal system.')
], 'PROJECT VISION')

# 3. Slide 3: Project Background & Problem Statement
add_content_slide(prs, 'Problem Statement & Existing Gaps', [
    ('Poor Planning & Scope Creep', 'Many academic projects fail due to unrealistic timelines, poor technological choices, and underspecified initial scope boundaries.'),
    ('Faculty Bandwidth Constraints', 'Professors frequently supervise 10-20 student teams simultaneously, limiting personalized guidance available to each team when stuck.'),
    ('Lack of Automated Project Blueprints', 'Students struggle to convert abstract project ideas into actionable, week-by-week implementation roadmaps.'),
    ('Static & Manual Documentation', 'Manual preparation of project synopses, progress reports, and defense slide decks consumes significant student development time.')
], 'BACKGROUND & CHALLENGES')

# 4. Slide 4: Team Role Allocation Table
add_table_slide(prs, 'Team Member Role & Responsibility Matrix', 
    ['Team Member', 'Project Role', 'Assignment Level', 'Core Responsibilities & Modules'],
    [
        ['Harshitha H S', 'Project Lead & System Architect', 'MAJOR ROLE', 'Multi-Agent Pipeline Architecture, Feasibility Agent, Scope Agent, Tech Stack Agent, Milestone Agent & Integration Lead.'],
        ['Vamshi Krishna', 'Backend Database & API Specialist', 'Minor Role', 'MySQL ai_project_mentor Schema, DB Connection Pooling (db.js), REST API Routing & Google/GitHub OAuth Sync.'],
        ['Muthumenen M', 'Frontend UI/UX Developer', 'Minor Role', 'Student Profile Onboarding (register.html), Skill Assessment Matrix (assessment.html), & Control Hub UI.'],
        ['Lohith Raj', 'Risk & AI Mentor Specialist', 'Minor Role', 'Risk Assessment & Mitigation Agent, Nova AI Mentor Chat Desk (ai-mentor.html), & Local Classifier Fallbacks.'],
        ['Harsh', 'Report Exporter & Faculty View Developer', 'Minor Role', 'On-Demand Document Exporter (report-view.html), Faculty Dashboard (faculty.html), & Teacher Feedback Modal System.']
    ], 'TEAM MANAGEMENT')

# 5. Slide 5: System Architecture & Workflow Pipeline
add_content_slide(prs, 'Multi-Agent Architecture & Pipeline Workflow', [
    ('Step 1: Student Input Submission', 'Student registers profile, completes skill assessment (0-100 score), and submits a 2-3 line project idea in the submission desk.'),
    ('Step 2: Unified Multi-Agent Controller Hub', 'Express API triggers the unified agent controller (/api/projects/agent/unified) orchestrating 5 autonomous agent workers.'),
    ('Step 3: Sequential Agent Processing', 'Feasibility Agent (Viability Index 88%) -> Scope Agent (MVP Boundaries) -> Tech Stack Agent (React/Express/MySQL) -> Timeline Agent (~80h Budget).'),
    ('Step 4: Output Rendering & Persistence', 'Agent outputs stream in real-time to the execution log terminal window and persist in localStorage / MySQL DB for live rendering.')
], 'SYSTEM ARCHITECTURE')

# 6. Slide 6: Module 1 — Student Profile & Skill Assessment
add_content_slide(prs, 'Module 1: Student Profile & Skill Assessment Matrix', [
    ('Competency Score Calculation', 'Calculates student technical competency out of 100 based on self-assessed programming languages, frameworks, and tools.'),
    ('Personalized Guidance Seeding', 'Feeds competency metrics into downstream agents to tailor technology stack recommendations and milestone time estimates.'),
    ('Interactive Onboarding UI', 'Responsive registration form built in register.html & assessment.html with visual rating sliders and skill badges.')
], 'CORE MODULES')

# 7. Slide 7: Module 2 — Project Submission & Feasibility Agent
add_content_slide(prs, 'Module 2: Project Idea Submission & Feasibility Analysis Agent', [
    ('Pre-Loaded Academic Template Cards', 'Provides 4 quick-start academic project templates (Healthcare AI, E-Commerce, IoT Tracker, Agent Mentor) for single-click filling.'),
    ('Feasibility Viability Index', 'Evaluates technical feasibility, resource requirements, and academic complexity, returning an 88% Viability Score.'),
    ('Automatic Pipeline Trigger', 'Submitting an idea automatically redirects the student to the Agent Hub and launches all 5 AI agents instantly.')
], 'CORE MODULES')

# 8. Slide 8: Module 3 — Scope Definition & Tech Stack Agent
add_content_slide(prs, 'Module 3: Scope Definition & Tech Stack Recommendation Agent', [
    ('MVP Scope Specification', 'Clearly demarcates core Must-Have MVP features from out-of-scope secondary enhancements to prevent scope creep.'),
    ('Intelligent Tech Stack Reasoning', 'Recommends modern frontend (HTML5/CSS3/Vanilla JS), backend (Node.js Express), and database (MySQL) with clear trade-off justifications.'),
    ('Architectural Guidelines', 'Provides database table structure schemas and REST API endpoint blueprints for rapid student setup.')
], 'CORE MODULES')

# 9. Slide 9: Module 4 — Milestone & Resource Timeline Agent
add_content_slide(prs, 'Module 4: Milestone & Resource Timeline Planning Agent', [
    ('Time-Budgeted Effort Calculator', 'Computes estimated development effort (~80 total student hours) tailored to academic project course durations.'),
    ('Week-by-Week Execution Matrix', 'Generates a structured 4-Milestone execution plan with clear task dependencies and deliverable checkpoints.'),
    ('Dynamic Check-in Alignment', 'Feeds milestone targets into the Weekly Check-in Desk for ongoing progress tracking.')
], 'CORE MODULES')

# 10. Slide 10: Module 5 — Risk Assessment & Mitigation Agent
add_content_slide(prs, 'Module 5: Risk Assessment & Mitigation Agent Engine', [
    ('Automated Risk Identification', 'Detects potential execution bottlenecks, API rate limits, model inference latency, and data dependency risks.'),
    ('Plan-B Resolution Strategies', 'Provides actionable fallback solutions (e.g., lightweight local DistilBERT classifiers when API latency exceeds 1.5s).'),
    ('Proactive Warning Badges', 'Renders visual risk callouts in the Agent Hub to alert students before starting development.')
], 'CORE MODULES')

# 11. Slide 11: Module 6 — Conversational AI Mentoring (Nova AI)
add_content_slide(prs, 'Module 6: Conversational Mentor Interaction Desk', [
    ('Nova AI Interactive Chat Desk', 'Provides 24/7 conversational mentoring for debugging code, clarifying architectural queries, and explaining concepts.'),
    ('Quick Action Prompt Chips', 'Pre-configured prompt shortcuts for instant answers on MySQL setup, REST APIs, and unit test verification.'),
    ('Context-Aware Guidance', 'Maintains active project details in session context to deliver highly relevant, targeted advice.')
], 'CORE MODULES')

# 12. Slide 12: Module 7 — On-Demand Document Exporter
add_content_slide(prs, 'Module 7: On-Demand Academic Document Generation Engine', [
    ('Academic Synopsis Exporter', 'Generates structured project synopses with problem statements, objectives, and module breakdowns.'),
    ('Methodology & Progress Reports', 'Auto-compiles technical methodologies, weekly progress summaries, and feature check-in logs.'),
    ('9-Slide Defense Presentation Outline', 'Renders a complete 9-slide PowerPoint outline with print-to-PDF formatting for academic project defense.')
], 'CORE MODULES')

# 13. Slide 13: Module 8 — Faculty Monitoring Dashboard
add_content_slide(prs, 'Module 8: Faculty Monitoring Dashboard System', [
    ('Centralized Student Team Overview', 'Provides faculty members with a single dashboard to passively monitor all assigned student project teams.'),
    ('Student Team Health Indicators', 'Displays visual status tags (+ On Track [Green], Needs Guidance [Yellow], At Risk [Red]).'),
    ('Auto-Generated Summary Cards', 'Displays AI-compiled project progress summaries without requiring faculty to manually inspect code.')
], 'CORE MODULES')

# 14. Slide 14: Interactive Faculty Feedback System
add_content_slide(prs, 'Interactive Faculty Guidance & Feedback System', [
    ('Faculty Guidance Feedback Modal', 'Enables teachers to select a student team, choose an evaluation tag, and type custom technical guidance notes.'),
    ('Live Student Banner Notification', 'Submitting feedback instantly triggers a prominent callout banner on student views ("Faculty Advisor Guidance Received").'),
    ('Persistent Guidance History', 'Saves teacher guidance notes in localStorage & MySQL database for continuous student reference.')
], 'FACULTY INTERACTION')

# 15. Slide 15: Security & Social Authentication
add_content_slide(prs, 'Google & GitHub Social Authentication Integration', [
    ('Single-Click OAuth Login', 'Integrated handleSocialAuth(provider) for Google, GitHub, and Microsoft single-click authorization.'),
    ('Profile Sync & Context Persistence', 'Automatically authenticates user profiles, sets localStorage context, and redirects to dashboard.html.'),
    ('Secure Session Management', 'Prevents unauthorized access while maintaining seamless single-sign-on across all platform pages.')
], 'SECURITY & AUTH')

# 16. Slide 16: Program Schedule & Timeline Overview
add_table_slide(prs, 'Official Program Schedule Overview (June 29 – August 20, 2026)',
    ['Milestone', 'Date Range', 'Target Focus & Deliverables', 'Active Team Members'],
    [
        ['Milestone 1 (Sprint 1)', 'June 29 – July 10, 2026', 'Multi-Agent System Architecture, MySQL DB Setup, Onboarding UI & Idea Submission', 'Harshitha H S (MAJOR), Vamshi, Muthumenen, Lohith, Harsh'],
        ['Milestone 2 (Sprint 2)', 'July 13 – July 24, 2026', 'Feasibility Agent, Scope Agent, Tech Stack Agent, Milestone Agent & Control Hub UI', 'Harshitha H S (MAJOR), Vamshi, Muthumenen, Lohith, Harsh'],
        ['Milestone 3 (Sprint 3)', 'July 27 – August 07, 2026', 'Weekly Check-in Engine, Risk Agent, Nova AI Chat Desk & Document Exporter', 'Harshitha H S (MAJOR), Vamshi, Muthumenen, Lohith, Harsh'],
        ['Milestone 4 (Sprint 4)', 'August 10 – August 20, 2026', 'Faculty Monitoring Dashboard, Feedback Modal, OAuth Login & Final Defense Prep', 'Harshitha H S (MAJOR), Vamshi, Muthumenen, Lohith, Harsh']
    ], 'PROGRAM TIMELINE')

# 17. Slide 17: Milestone 1 Sprint Task Division Table
add_table_slide(prs, 'Milestone 1 Task Division (June 29 – July 10, 2026)',
    ['US ID', 'Task Description', 'Start Date', 'End Date', 'Assignee', 'Status'],
    [
        ['US-001', 'Lead Multi-Agent Architecture & System Data Flow Design', '2026-06-29', '2026-07-10', 'Harshitha H S (MAJOR ROLE)', '3- Completed'],
        ['US-002', 'Design MySQL ai_project_mentor schema & users table', '2026-06-29', '2026-07-03', 'Vamshi Krishna', '3- Completed'],
        ['US-003', 'Develop Student Registration & Profile Onboarding Interface', '2026-06-30', '2026-07-06', 'Muthumenen M', '3- Completed'],
        ['US-004', 'Build Competency Skill Assessment Matrix UI & Scoring Logic', '2026-07-02', '2026-07-08', 'Lohith Raj', '3- Completed'],
        ['US-005', 'Develop Project Submission Page & Pre-Loaded Templates', '2026-07-06', '2026-07-10', 'Harsh', '3- Completed']
    ], 'AGILE SPRINT 1')

# 18. Slide 18: Milestone 2 Sprint Task Division Table
add_table_slide(prs, 'Milestone 2 Task Division (July 13 – July 24, 2026)',
    ['US ID', 'Task Description', 'Start Date', 'End Date', 'Assignee', 'Status'],
    [
        ['US-006', 'Develop Feasibility Analysis Agent Engine (88% Viability Score)', '2026-07-13', '2026-07-17', 'Harshitha H S (MAJOR ROLE)', '3- Completed'],
        ['US-007', 'Develop Scope Agent & Tech Stack Recommendation Agent', '2026-07-17', '2026-07-24', 'Harshitha H S (MAJOR ROLE)', '3- Completed'],
        ['US-008', 'Develop Express REST API Endpoints (/api/projects/agent/unified)', '2026-07-13', '2026-07-24', 'Vamshi Krishna', '3- Completed'],
        ['US-009', 'Build Multi-Agent Control Hub UI & Pipeline Trigger Cards', '2026-07-15', '2026-07-21', 'Muthumenen M', '3- Completed'],
        ['US-010', 'Implement Real-Time Execution Log Stream Terminal Window', '2026-07-17', '2026-07-23', 'Lohith Raj', '3- Completed'],
        ['US-011', 'Build Active Project Summary Display Header & Storage Payload', '2026-07-20', '2026-07-24', 'Harsh', '3- Completed']
    ], 'AGILE SPRINT 2')

# 19. Slide 19: Milestone 3 Sprint Task Division Table
add_table_slide(prs, 'Milestone 3 Task Division (July 27 – August 07, 2026)',
    ['US ID', 'Task Description', 'Start Date', 'End Date', 'Assignee', 'Status'],
    [
        ['US-012', 'Build Dynamic Weekly Check-in Progress Engine & Status Badge Shifter', '2026-07-27', '2026-08-04', 'Harshitha H S (MAJOR ROLE)', '3- Completed'],
        ['US-013', 'Implement Backend Chat Endpoints & Project Context Storage', '2026-07-27', '2026-08-04', 'Vamshi Krishna', '3- Completed'],
        ['US-014', 'Build Nova AI Mentor Chat UI & Quick Action Prompt Chips', '2026-07-29', '2026-08-06', 'Muthumenen M', '3- Completed'],
        ['US-015', 'Develop Risk Assessment Agent & Plan-B Mitigation Fallback Engines', '2026-07-27', '2026-07-31', 'Lohith Raj', '3- Completed'],
        ['US-016', 'Develop On-Demand Document Exporter (Synopsis, Slides PDF)', '2026-08-03', '2026-08-07', 'Harsh', '3- Completed']
    ], 'AGILE SPRINT 3')

# 20. Slide 20: Milestone 4 Sprint Task Division Table
add_table_slide(prs, 'Milestone 4 Task Division (August 10 – August 20, 2026)',
    ['US ID', 'Task Description', 'Start Date', 'End Date', 'Assignee', 'Status'],
    [
        ['US-017', 'Lead System Integration, Prompt Tuning & Final Presentation Prep', '2026-08-10', '2026-08-20', 'Harshitha H S (MAJOR ROLE)', '3- Completed'],
        ['US-018', 'Implement Google & GitHub OAuth Single-Click Authentication Sync', '2026-08-12', '2026-08-19', 'Vamshi Krishna', '3- Completed'],
        ['US-019', 'Reconfigure Student Command Dashboard with KPI Cards & Launcher Grid', '2026-08-10', '2026-08-14', 'Muthumenen M', '3- Completed'],
        ['US-020', 'Conduct End-to-End System Testing & Unit Test Case Verification', '2026-08-14', '2026-08-20', 'Lohith Raj', '3- Completed'],
        ['US-021', 'Develop Faculty Monitoring Dashboard & Feedback Modal System', '2026-08-11', '2026-08-18', 'Harsh', '3- Completed']
    ], 'AGILE SPRINT 4')

# 21. Slide 21: Daily Standups & Problem Resolutions Table
add_table_slide(prs, 'Daily Standup Impediments & Resolution Log',
    ['Sprint / Date', 'Impediments Encountered', 'Technical Action Taken'],
    [
        ['Sprint 1 (2026-07-02)', 'Need structured DB schema for student profiles & project metrics.', 'Vamshi Krishna created MySQL users table; Harshitha H S reviewed data models.'],
        ['Sprint 1 (2026-07-07)', 'Template cards not populating submission text area upon click.', 'Muthumenen M & Harsh updated click event handlers in submit-project.js.'],
        ['Sprint 2 (2026-07-16)', 'Gemini API key rate limits during initial agent pipeline testing.', 'Harshitha H S implemented rich offline fallback generators in aiAgents.js.'],
        ['Sprint 3 (2026-07-30)', 'Real-time sentiment inference latency exceeding 1.5s.', 'Lohith Raj & Harshitha H S deployed lightweight local DistilBERT classifiers (<500ms).'],
        ['Sprint 4 (2026-08-18)', 'Faculty needed direct method to send feedback when teams get stuck.', 'Harsh & Harshitha H S built Faculty Feedback Modal & live banner notification system.']
    ], 'STANDUP MEETINGS')

# 22. Slide 22: Unit Test Verification Table
add_table_slide(prs, 'Unit Test Verification Results Table',
    ['Test ID', 'Test Case Name', 'Tested Condition', 'Expected vs Actual Result', 'Status'],
    [
        ['UT-001', 'Student Profile Registration', 'Form validation & skill scoring', 'Saved in MySQL DB & competency score calculated', 'PASS'],
        ['UT-002', 'Project Idea Submission', 'Click template card & press Enter', 'Saves idea state & redirects seamlessly to Agent Hub', 'PASS'],
        ['UT-003', 'Multi-Agent Blueprint Execution', 'Click Run Complete Pipeline', 'Generates Feasibility (88%), Scope, Stack & Timeline', 'PASS'],
        ['UT-004', 'Weekly Check-in Badge Shift', 'Submit Week 1 check-in', 'Milestone 1 turns Green Completed (100%), M2 turns Active', 'PASS'],
        ['UT-005', 'Document Exporter PDF Generation', 'Click Academic Synopsis & Print', 'Renders project synopsis & opens PDF print dialog', 'PASS'],
        ['UT-006', 'Faculty Feedback Modal & Banner', 'Faculty sends guidance note', 'Displays Faculty Guidance Received banner on student dashboard', 'PASS'],
        ['UT-007', 'Google & GitHub OAuth Login', 'Click Continue with Google', 'Sets authenticated user session & redirects to dashboard.html', 'PASS']
    ], 'TESTING & QA')

# 23. Slide 23: Defect Tracker Resolution Table
add_table_slide(prs, 'Defect Tracker Log & QA Summary Table',
    ['ID', 'Submitted By', 'Description of Defect', 'Assigned To', 'Type', 'Action Taken & Verification Status'],
    [
        ['1', 'Muthumenen M', 'Template cards not filling text area upon click', 'Muthumenen M', 'UI', 'Fixed click event listener in submit-project.js. (Closed)'],
        ['2', 'Harshitha H S', 'API rate limit error during agent blueprint generation', 'Harshitha H S (MAJOR)', 'Logical', 'Implemented rich offline fallback generators in aiAgents.js. (Closed)'],
        ['3', 'Lohith Raj', 'High inference latency (>1.5s) during sentiment analysis', 'Lohith Raj', 'Performance', 'Switched to lightweight local DistilBERT classifier (<500ms). (Closed)'],
        ['4', 'Harshitha H S', 'Milestone badges showing 100% completed prematurely', 'Harshitha H S (MAJOR)', 'Logical', 'Added completedCheckinWeek localStorage tracking (0% initial state). (Closed)'],
        ['5', 'Harsh', 'Excessive white space on command dashboard layout', 'Harshitha H S (MAJOR)', 'UI', 'Reconfigured dashboard.html with 4 KPI cards & launcher grid. (Closed)'],
        ['6', 'Harsh', 'Faculty unable to send guidance note to stuck teams', 'Harsh', 'Logical', 'Built Interactive Faculty Feedback Modal & live banner system. (Closed)'],
        ['7', 'Vamshi Krishna', 'Social login buttons not creating authenticated context', 'Vamshi Krishna', 'Standards', 'Integrated handleSocialAuth handler for Google/GitHub OAuth. (Closed)']
    ], 'DEFECT TRACKER')

# 24. Slide 24: Evaluation Criteria Mapping
add_content_slide(prs, 'Evaluation Criteria Mapping & Achievements', [
    ('1. Quality of Automated Blueprint', 'EXCEEDED — Single-click input generates complete Feasibility (88%), Scope MVP, Tech Stack, and Milestone Plan automatically.'),
    ('2. Effectiveness of Individual Agents', 'EXCEEDED — 5 specialized AI agents work in unison with rich offline fallback generators to ensure 100% uptime.'),
    ('3. Conversational Mentoring Relevance', 'EXCEEDED — Nova AI Chat Desk provides 24/7 contextual guidance with pre-configured quick prompt chips for instant answers.'),
    ('4. Documentation & Faculty Dashboard Quality', 'EXCEEDED — On-Demand Exporter renders PDF synopses; Faculty Dashboard features live guidance modals & banner notifications.')
], 'EVALUATION & RESULTS')

# 25. Slide 25: Project Conclusion & Q&A
add_content_slide(prs, 'Project Conclusion, Future Scope & Q&A', [
    ('Project Conclusion', 'The AI Academic Project Mentor successfully bridges the gap between students and mentors, converting 2-line ideas into actionable blueprints and empowering faculty supervision.'),
    ('Future Scope & Scalability', 'Expansion to integrate directly with University LMS platforms (Canvas, Moodle), automated Git commit tracking, and AI code review capabilities.'),
    ('Thank You & Questions', 'Thank you! We welcome any questions, feedback, and technical discussion.')
], 'CONCLUSION & Q&A')

# Save PPTX to Downloads and NEW SPRING
out_path_1 = r'C:\Users\Harshita\Downloads\NEW SPRING\AI_Academic_Project_Mentor_25_Slide_Presentation.pptx'
out_path_2 = r'C:\Users\Harshita\Downloads\AI_Academic_Project_Mentor_25_Slide_Presentation.pptx'

prs.save(out_path_1)
prs.save(out_path_2)

print('SUCCESSFULLY GENERATED 25-SLIDE PRESENTATION AT:')
print(out_path_1)
print(out_path_2)
